from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

# Basic styling function
def style_title(shape):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = 'Arial'
            run.font.color.rgb = RGBColor(0, 128, 128) # Teal
            run.font.bold = True

def add_slide(prs, title, content_lines, image_path=None):
    if image_path and os.path.exists(image_path):
        slide_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        style_title(title_shape)
        
        # Add image
        try:
            slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
        except Exception as e:
            print(f"Failed to add image {image_path}: {e}")
    else:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        style_title(title_shape)
        
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        for i, line in enumerate(content_lines):
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "AquaSense"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204) # Blue

subtitle.text = ("Smart Water Monitoring System\n\n"
                 "Team Members: Krish Mahajan, Shreyash Chougule, Pranav Chougule, Varad Dake\n"
                 "D.K.T.E. Society’s Textile and Engineering Institute, Ichalkaranji")

# Slide 2: Problem Statement
add_slide(prs, "Problem Statement", [
    "• Significant water wastage in shared environments.",
    "• Lack of real-time water consumption monitoring.",
    "• Absence of section-wise or area-wise tracking.",
    "• Frequent tank overflows and unexpected shortages."
])

# Slide 3: Aim of the Project
add_slide(prs, "Aim of the Project", [
    "• To develop an IoT-based smart water monitoring and management system.",
    "• Provide real-time, multi-area water intelligence to prevent wastage.",
    "• Facilitate sustainable water management practices."
])

# Slide 4: Objectives
add_slide(prs, "Objectives", [
    "• Deploy flow sensors to monitor area-wise consumption independently.",
    "• Analyze real-time flow and tank level data.",
    "• Regulate distribution using servo-controlled valves.",
    "• Provide actionable insights through a live dashboard."
])

# Slide 5: Scope and Limitations
add_slide(prs, "Scope and Limitations", [
    "Scope:",
    "• Hostels, apartments, institutions, and agriculture.",
    "Limitations:",
    "• Dependent on continuous power and Wi-Fi connectivity.",
    "• Hardware maintenance required periodically."
])

# Slide 6: Existing vs Proposed System
add_slide(prs, "Existing System vs Proposed System", [
    "Existing System:",
    "• Manual monitoring, no section-wise data.",
    "• Delayed overflow detection.",
    "",
    "Proposed System (AquaSense):",
    "• Automated IoT monitoring.",
    "• Area-wise analytics and instant anomaly alerts."
])

# Slide 7: Proposed Solution (Overview)
add_slide(prs, "Proposed Solution (AquaSense Overview)", [
    "• Integrated system using ESP32/ESP8266 microcontroller.",
    "• Collects data from flow and tank sensors.",
    "• Processes and displays data on web/mobile dashboard.",
    "• Includes manual and Auto Mode valve regulation.",
    "• Cloud integration for persistent storage and remote access."
])

# Slide 8: System Architecture Diagram
add_slide(prs, "System Architecture Diagram", [], "Diagrams/architecture dia.png")

# Slide 9: Working Flow
add_slide(prs, "Working Flow", [
    "• Step 1: Sensors detect flow and tank levels.",
    "• Step 2: Microcontroller processes signals.",
    "• Step 3: Data transmitted to backend via Wi-Fi.",
    "• Step 4: System analyzes consumption and detects anomalies.",
    "• Step 5: Dashboard updates & Telegram alerts sent."
])

# Slide 10: Modules Explanation
add_slide(prs, "Modules Explanation", [
    "• Sensor Simulation: Simulates realistic flow data.",
    "• Backend Processing: Handles data aggregation and logic.",
    "• Adaptive Aeration Logic: Optimizes water flow.",
    "• Anomaly Detection: Identifies leaks and overflows.",
    "• Dashboard: Real-time user interface for monitoring.",
    "• Telegram Integration: Instant notifications and control."
])

# Slide 11: Algorithm Logic
add_slide(prs, "Algorithm Logic", [
    "1. Initialize sensors and connect to network.",
    "2. Read data continuously.",
    "3. If usage exceeds threshold, trigger anomaly.",
    "4. If tank full/empty, actuate valve.",
    "5. Log data and notify user."
])

# Slide 12: Use Case Diagram
add_slide(prs, "Use Case Diagram", [], "Diagrams/usecase.jpeg")

# Slide 13: Class Diagram
add_slide(prs, "Class Diagram", [], "Diagrams/class diagram.jpeg")

# Slide 14: Sequence Diagram
add_slide(prs, "Sequence Diagram", [], "Diagrams/Sequence.png")

# Slide 15: Activity Diagram
add_slide(prs, "Activity Diagram", [], "Diagrams/Activity diagram.png")

# Slide 16: Implementation Details
add_slide(prs, "Implementation Details (Tools & Technologies)", [
    "• Hardware: ESP32/ESP8266, Flow Sensors, Servo Valves.",
    "• Software: Python, Flask backend, Web/Mobile Dashboard.",
    "• Cloud Storage: Firebase / Firestore.",
    "• Communication: Wi-Fi, Telegram Bot API."
])

# Slide 17: Results & Performance
add_slide(prs, "Results & Performance", [
    "• Achieved 22–35% water saving.",
    "• Fast 1.5 sec anomaly detection.",
    "• Accurate area-wise consumption tracking.",
    "• Reliable automated valve control."
])

# Slide 18: Outputs / Screens
add_slide(prs, "Outputs / Screens", [
    "• Dashboard: Real-time metrics overview.",
    "• Analytics: Historical consumption graphs.",
    "• Control Panel: Manual valve override.",
    "• Telegram: Alerts on leakages or low tank levels."
])

# Slide 19: Testing
add_slide(prs, "Testing", [
    "• Integration Testing: Seamless hardware-software communication.",
    "• System Testing: Validated end-to-end functionality under stress.",
    "• Robustness verified for continuous 24/7 operation."
])

# Slide 20: Future Scope
add_slide(prs, "Future Scope", [
    "• Integration with AI for predictive maintenance.",
    "• Solar-powered sensor nodes for sustainability.",
    "• Expansion to smart city infrastructure.",
    "• App-based notifications for granular user control."
])

# Slide 21: Applications
add_slide(prs, "Applications", [
    "• Residential complexes and hostels.",
    "• Educational institutions and hospitals.",
    "• Precision agriculture and plantation.",
    "• Industrial water management systems."
])

# Slide 22: Conclusion
add_slide(prs, "Conclusion", [
    "• AquaSense provides an effective, low-cost solution for water management.",
    "• Real-time tracking and automated controls significantly reduce wastage.",
    "• Simple, user-friendly interface enhances system usability."
])

# Slide 23: Thank You
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
subtitle.text = "Any Questions?"

prs.save('AquaSense_Presentation.pptx')
print("Presentation generated successfully: AquaSense_Presentation.pptx")
